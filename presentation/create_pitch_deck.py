from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Load color palette
main_color = RGBColor(255, 99, 71)      # Tomato Red
accent_color = RGBColor(255, 215, 0)    # Gold
secondary_color = RGBColor(70, 130, 180) # Steel Blue
bg_color = RGBColor(255, 255, 255)      # White
text_color = RGBColor(70, 130, 180)     # Steel Blue

# Slide content outline
slides = [
    ("Introduction", ["Introducing Mr. Gift: Your Ultimate Gift Ordering App", "Welcome and introduction", "Overview of the app's purpose and benefits"]),
    ("The Gift-Giving Dilemma", ["Challenges in finding the perfect gift for weddings and birthdays", "Time-consuming, overwhelming, uncertainty", "Need for a solution"]),
    ("Meet Mr. Gift", ["Innovative solution to the gift-giving dilemma", "AI-powered, user-friendly, personalized recommendations"]),
    ("Key Features of Mr. Gift", ["User registration and personalized profiles", "Extensive gift catalogue for weddings and birthdays", "AI-powered recommendation system", "Secure payment integration", "Order tracking and delivery updates"]),
    ("How Mr. Gift Works", ["Step-by-step user journey:", "1. Registration and profile creation", "2. Browse and select gifts", "3. AI-driven recommendations", "4. Secure payment", "5. Order tracking"]),
    ("AI Recommendation System", ["How AI analyzes preferences, occasions, and history", "Examples of personalized recommendations"]),
    ("Market Opportunity", ["Market size and trends", "Demand for convenience and personalization"]),
    ("Revenue Model", ["Commission on sales", "Premium features/subscription", "Advertising and promotions", "Revenue growth projections"]),
    ("Marketing and Growth Strategy", ["Digital marketing, social media, influencer partnerships", "Event planning partnerships", "Referral program", "User feedback and continuous improvement"]),
    ("Competitive Advantage", ["AI recommendation system", "User-friendly interface", "Convenience and time-saving", "Personalized experience"]),
    ("Budget", ["App development (front-end, back-end, QA, PM)", "Hosting and infrastructure", "Third-party API integration (AI, payment)", "Marketing and user acquisition", "Ongoing maintenance and support"]),
    ("Investment Opportunity", ["Investment required and equity offered", "ROI potential"]),
    ("Conclusion and Contact", ["Recap key points and benefits", "Call to action for investors", "Contact information"]),
    ("Q&A", ["Open for questions"]),
]

prs = Presentation()

# Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(0.5)
top = Inches(0.5)
logo_path = "branding/mr_gift_logo.png"
slide.shapes.add_picture(logo_path, left, top, height=Inches(1.5))
title_box = slide.shapes.add_textbox(Inches(2.5), Inches(0.5), Inches(6), Inches(2))
tf = title_box.text_frame
tf.text = "Introducing Mr. Gift\nYour Ultimate Gift Ordering App"
tf.paragraphs[0].font.size = Pt(40)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = main_color

# Add slides
for title, bullets in slides[1:]:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = main_color
    # Bullets
    left = Inches(0.8)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        p.level = 0
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = text_color

# Add color palette slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Brand Color Palette"
slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
slide.shapes.title.text_frame.paragraphs[0].font.bold = True
slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = main_color
palette_img = "branding/mr_gift_palette.png"
slide.shapes.add_picture(palette_img, Inches(1), Inches(2), height=Inches(1))

prs.save("presentation/mr_gift_pitch_deck.pptx")
print("Presentation created: presentation/mr_gift_pitch_deck.pptx")